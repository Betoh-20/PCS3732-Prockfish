#!/usr/bin/env python3
"""Gera a apresentação final (semana 4) em docs/Apresentacao-Final.pptx.

O conteúdo vem dos relatórios das semanas 1-3 (docs/Relatorio-*.{md,tex}),
do README e de arquitetura_projeto.md; as figuras vêm de docs/diagramas/.

    pip install python-pptx && python docs/gerar_apresentacao.py
"""

from pathlib import Path

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Emu, Inches, Pt

RAIZ = Path(__file__).resolve().parent.parent
DIAGRAMAS = RAIZ / "docs" / "diagramas"
SAIDA = RAIZ / "docs" / "Apresentacao-Final.pptx"

# Paleta: fundo escuro dos próprios diagramas Mermaid + âmbar de destaque.
FUNDO = RGBColor(0x0F, 0x1B, 0x2D)
TITULO = RGBColor(0xF2, 0xF4, 0xF8)
TEXTO = RGBColor(0xD6, 0xDE, 0xE8)
FRACO = RGBColor(0x9F, 0xB3, 0xC8)
ACENTO = RGBColor(0xE2, 0xB7, 0x14)
VERDE = RGBColor(0x4C, 0xC3, 0x8A)
VERMELHO = RGBColor(0xE9, 0x45, 0x60)
CARTAO = RGBColor(0x17, 0x27, 0x40)
BORDA = RGBColor(0x2A, 0x40, 0x63)

L, A = Inches(13.333), Inches(7.5)  # 16:9
MARGEM = Inches(0.75)


# --------------------------------------------------------------------------
# Blocos de montagem
# --------------------------------------------------------------------------
def slide_vazio(prs):
    s = prs.slides.add_slide(prs.slide_layouts[6])
    fundo = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, L, A)
    fundo.fill.solid()
    fundo.fill.fore_color.rgb = FUNDO
    fundo.line.fill.background()
    fundo.shadow.inherit = False
    return s


def caixa(s, x, y, cx, cy):
    tb = s.shapes.add_textbox(x, y, cx, cy)
    tf = tb.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    return tf


def par(tf, texto, tam=18, cor=TEXTO, negrito=False, espaco=8,
        marcador=None, primeiro=False, alinhamento=PP_ALIGN.LEFT):
    p = tf.paragraphs[0] if primeiro else tf.add_paragraph()
    p.alignment = alinhamento
    p.space_after = Pt(espaco)
    if marcador:
        r = p.add_run()
        r.text = marcador + "  "
        r.font.size = Pt(tam)
        r.font.color.rgb = ACENTO
        r.font.bold = True
    r = p.add_run()
    r.text = texto
    r.font.size = Pt(tam)
    r.font.bold = negrito
    r.font.color.rgb = cor
    return p


def cabecalho(s, titulo, subtitulo=None):
    tf = caixa(s, MARGEM, Inches(0.45), L - 2 * MARGEM, Inches(1.0))
    par(tf, titulo, tam=32, cor=TITULO, negrito=True, espaco=2, primeiro=True)
    if subtitulo:
        par(tf, subtitulo, tam=15, cor=FRACO, espaco=0)
    linha = s.shapes.add_shape(
        MSO_SHAPE.RECTANGLE, MARGEM, Inches(1.42), Inches(1.5), Pt(3))
    linha.fill.solid()
    linha.fill.fore_color.rgb = ACENTO
    linha.line.fill.background()
    linha.shadow.inherit = False


def rodape(s, texto):
    tf = caixa(s, MARGEM, A - Inches(0.62), L - 2 * MARGEM, Inches(0.35))
    par(tf, texto, tam=11, cor=FRACO, espaco=0, primeiro=True)


def cartao(s, x, y, cx, cy, cor_borda=BORDA, preenchimento=CARTAO):
    f = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, x, y, cx, cy)
    f.adjustments[0] = 0.06
    f.fill.solid()
    f.fill.fore_color.rgb = preenchimento
    f.line.color.rgb = cor_borda
    f.line.width = Pt(1.25)
    f.shadow.inherit = False
    f.text_frame.word_wrap = True
    return f


def cartao_texto(s, x, y, cx, cy, titulo, linhas, cor=ACENTO):
    f = cartao(s, x, y, cx, cy, cor_borda=cor)
    tf = f.text_frame
    tf.margin_left = tf.margin_right = Inches(0.22)
    tf.margin_top = tf.margin_bottom = Inches(0.18)
    tf.vertical_anchor = MSO_ANCHOR.TOP
    par(tf, titulo, tam=17, cor=cor, negrito=True, espaco=6, primeiro=True)
    for linha in linhas:
        par(tf, linha, tam=13, cor=TEXTO, espaco=4)
    return f


def figura(s, nome, topo, altura=None, largura=None):
    """Insere um diagrama centralizado horizontalmente."""
    caminho = str(DIAGRAMAS / nome)
    if altura:
        pic = s.shapes.add_picture(caminho, 0, topo, height=altura)
    else:
        pic = s.shapes.add_picture(caminho, 0, topo, width=largura)
    pic.left = int((L - pic.width) / 2)
    return pic


def mono(s, x, y, cx, cy, linhas, tam=13):
    f = cartao(s, x, y, cx, cy, cor_borda=BORDA, preenchimento=RGBColor(0x0A, 0x12, 0x1F))
    tf = f.text_frame
    tf.margin_left = tf.margin_right = Inches(0.2)
    tf.margin_top = tf.margin_bottom = Inches(0.14)
    for i, (texto, cor) in enumerate(linhas):
        p = par(tf, texto, tam=tam, cor=cor, espaco=2, primeiro=(i == 0))
        for r in p.runs:
            r.font.name = "Consolas"
    return f


# --------------------------------------------------------------------------
# Slides
# --------------------------------------------------------------------------
def s01_capa(prs):
    s = slide_vazio(prs)
    faixa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), A)
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = ACENTO
    faixa.line.fill.background()
    faixa.shadow.inherit = False

    tf = caixa(s, Inches(1.1), Inches(2.1), Inches(10.5), Inches(3.2))
    par(tf, "PCS3732 — Laboratório de Processadores", tam=16, cor=ACENTO,
        negrito=True, espaco=14, primeiro=True)
    par(tf, "Tabuleiro de Xadrez Eletrônico", tam=48, cor=TITULO,
        negrito=True, espaco=4)
    par(tf, "com integração a chess engine e a partidas online", tam=26,
        cor=TEXTO, espaco=22)
    par(tf, "Apresentação Final  ·  Grupo W  ·  Agosto de 2026", tam=15,
        cor=FRACO, espaco=0)

    # Miniatura de tabuleiro no canto direito.
    lado = Inches(0.34)
    x0, y0 = Inches(10.55), Inches(4.55)
    for i in range(8):
        for j in range(8):
            q = s.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, x0 + j * lado, y0 + i * lado, lado, lado)
            q.fill.solid()
            q.fill.fore_color.rgb = (RGBColor(0x1C, 0x2E, 0x4A)
                                     if (i + j) % 2 else RGBColor(0x30, 0x4A, 0x70))
            q.line.fill.background()
            q.shadow.inherit = False


def s02_motivacao(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Motivação", "Duas formas de jogar xadrez que não conversam entre si")

    larg = Inches(3.55)
    y = Inches(1.95)
    cartao_texto(s, MARGEM, y, larg, Inches(2.5), "Tabuleiro físico", [
        "Peças reais, presença, o gesto de mover a peça",
        "Sem oponente disponível, sem análise, sem registro",
    ], cor=FRACO)
    cartao_texto(s, MARGEM + larg + Inches(0.3), y, larg, Inches(2.5),
                 "Xadrez digital", [
        "Engine forte a qualquer hora, milhões de oponentes online",
        "A partida acontece numa tela — o tabuleiro some",
    ], cor=FRACO)
    cartao_texto(s, MARGEM + 2 * (larg + Inches(0.3)), y, larg, Inches(2.5),
                 "Nossa proposta", [
        "Peças físicas de verdade sobre um tabuleiro instrumentado",
        "Oponente (Stockfish ou Lichess) representado na GUI",
    ], cor=ACENTO)

    tf = caixa(s, MARGEM, Inches(5.0), L - 2 * MARGEM, Inches(1.6))
    par(tf, "O problema técnico", tam=20, cor=TITULO, negrito=True,
        espaco=10, primeiro=True)
    par(tf, "Detectar o movimento de peças físicas sem intervenção manual, "
            "traduzi-lo em jogada válida e devolver o estado completo da "
            "partida ao jogador — em tempo real, sobre hardware embarcado.",
        tam=17, cor=TEXTO, espaco=0)


def s03_objetivos(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Objetivos e requisitos", "Definidos na semana 1")

    larg = Inches(5.6)
    cartao_texto(s, MARGEM, Inches(1.95), larg, Inches(2.35),
                 "Requisitos funcionais", [
        "RF1 — Detectar as peças no tabuleiro físico e seus movimentos",
        "RF2 — Comunicar-se com uma chess engine como oponente",
        "RF3 — Exibir graficamente as duas partes do jogo (32 peças)",
    ])
    cartao_texto(s, MARGEM + larg + Inches(0.5), Inches(1.95), larg, Inches(2.35),
                 "Requisitos não funcionais", [
        "RNF1 — Detecção < 200 ms, com 100 % de precisão",
        "RNF2 — GUI sem latência perceptível (< 100 ms)",
        "RNF3 — Robustez contra falsos positivos (esbarrões)",
    ])

    tf = caixa(s, MARGEM, Inches(4.7), L - 2 * MARGEM, Inches(2.0))
    par(tf, "Objetivo geral", tam=20, cor=TITULO, negrito=True, espaco=10,
        primeiro=True)
    par(tf, "Projetar e implementar um tabuleiro de xadrez eletrônico que "
            "detecte movimentos de peças físicas, os comunique ao Stockfish "
            "ou à plataforma Lichess, e exiba o estado completo da partida "
            "em interface gráfica.", tam=17, cor=TEXTO, espaco=14)
    par(tf, "Metodologia: ciclo iterativo e incremental em entregas semanais, "
            "com branches integradas por pull request e validação em camadas "
            "(testes automatizados, mock de hardware, prova de conceito, "
            "sondas contra a API real).", tam=14, cor=FRACO, espaco=0)


def s04_solucao(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Solução proposta",
              "Tabuleiro instrumentado + Raspberry Pi + dois processos cooperando")

    itens = [
        ("Tabuleiro instrumentado",
         "Matriz 8×8 de reed switches com diodos anti-ghosting; ímãs de "
         "neodímio embutidos nas peças. 64 sensores em apenas 16 pinos GPIO."),
        ("Processo em C — baixo nível",
         "Varredura matricial, debouncing por amostragem consecutiva e "
         "detecção de diferenças. Emite eventos curtos de texto."),
        ("Processo em Python — aplicação",
         "Estado do jogo e validação com python-chess, interface UCI com o "
         "Stockfish, cliente da Lichess Board API e GUI em pygame."),
        ("Duas camadas de entrada intercambiáveis",
         "Reed switches ou teclado matricial 4×4 — mesmo protocolo IPC, "
         "aplicação idêntica nos dois casos."),
    ]
    y = Inches(1.95)
    for i, (titulo, corpo) in enumerate(itens):
        n = s.shapes.add_shape(MSO_SHAPE.OVAL, MARGEM, y, Inches(0.44), Inches(0.44))
        n.fill.solid()
        n.fill.fore_color.rgb = ACENTO
        n.line.fill.background()
        n.shadow.inherit = False
        ntf = n.text_frame
        ntf.margin_left = ntf.margin_right = ntf.margin_top = ntf.margin_bottom = 0
        par(ntf, str(i + 1), tam=15, cor=FUNDO, negrito=True, espaco=0,
            primeiro=True, alinhamento=PP_ALIGN.CENTER)

        tf = caixa(s, MARGEM + Inches(0.68), y - Inches(0.04),
                   L - 2 * MARGEM - Inches(0.68), Inches(1.1))
        par(tf, titulo, tam=18, cor=TITULO, negrito=True, espaco=3, primeiro=True)
        par(tf, corpo, tam=14, cor=TEXTO, espaco=0)
        y += Inches(1.22)


def s05_arquitetura(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Arquitetura do sistema",
              "Três camadas interligadas por comunicação entre processos (IPC)")

    larg = Inches(3.75)
    alt = Inches(3.35)
    y = Inches(1.95)
    espaco = Inches(0.42)
    x = MARGEM

    cartao_texto(s, x, y, larg, alt, "Hardware / sensoriamento", [
        "Peças com ímã de neodímio",
        "Matriz 8×8 de reed switches",
        "Diodos 1N4148 anti-ghosting",
        "GPIO do Raspberry Pi 3B+",
        "Alternativa: teclado matricial 4×4",
    ], cor=VERMELHO)

    x += larg + espaco
    cartao_texto(s, x, y, larg, alt, "Baixo nível — processo em C", [
        "main.c — CLI e escolha da camada",
        "reed_layer.c — varredura + debouncing",
        "keypad_layer.c — lances digitados",
        "board.c — espelho de ocupação",
        "ipc.c — serialização dos eventos",
        "lcd.c / gpio.c / runstate.c",
    ], cor=ACENTO)

    x += larg + espaco
    cartao_texto(s, x, y, larg, alt, "Aplicação — processo em Python", [
        "ipc_reader.py + move_interpreter.py",
        "game_state.py — python-chess",
        "stockfish_engine.py — protocolo UCI",
        "lichess_client.py — Board API",
        "gui.py — pygame",
        "menu.py / launcher.py — configuração",
    ], cor=VERDE)

    for cx in (MARGEM + larg + Inches(0.06), MARGEM + 2 * larg + espaco + Inches(0.06)):
        seta = s.shapes.add_shape(
            MSO_SHAPE.RIGHT_ARROW, cx, y + Inches(1.45), Inches(0.30), Inches(0.32))
        seta.fill.solid()
        seta.fill.fore_color.rgb = FRACO
        seta.line.fill.background()
        seta.shadow.inherit = False

    tf = caixa(s, MARGEM, Inches(5.55), Inches(4.4), Inches(0.9))
    par(tf, "GPIO", tam=12, cor=FRACO, espaco=0, primeiro=True,
        alinhamento=PP_ALIGN.CENTER)
    tf = caixa(s, MARGEM + larg + espaco, Inches(5.55), Inches(3.75), Inches(0.9))
    par(tf, "Named pipe / stdout", tam=12, cor=FRACO, espaco=0, primeiro=True,
        alinhamento=PP_ALIGN.CENTER)

    tf = caixa(s, MARGEM, Inches(5.95), L - 2 * MARGEM, Inches(0.9))
    par(tf, "Por que dois processos: o C dá controle direto de GPIO com "
            "temporização precisa; o Python traz o python-chess (validação e "
            "UCI). Named pipe é IPC nativo do Linux, suficiente para eventos "
            "curtos e esporádicos.", tam=13, cor=FRACO, espaco=0, primeiro=True)


def s06_fisica(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Arquitetura física", "Do ímã ao monitor")
    figura(s, "arq_fisica.png", Inches(2.35), largura=Inches(11.6))
    tf = caixa(s, MARGEM, Inches(5.25), L - 2 * MARGEM, Inches(1.6))
    par(tf, "Varredura matricial: 8 linhas em OUTPUT LOW, 8 colunas em "
            "INPUT_PULLUP, com um diodo 1N4148 em série por sensor. "
            "64 sensores lidos com 16 pinos.", tam=15, cor=TEXTO, espaco=8,
        primeiro=True)
    par(tf, "Prova de conceito validada em Arduino Mega 2560 e simulada no "
            "SimulIDE; versão final no Raspberry Pi 3B+, com acesso ao GPIO pelo "
            "processo em C (sem overhead de serial).", tam=15, cor=TEXTO, espaco=0)


def s07_ipc(prs):
    s = slide_vazio(prs)
    cabecalho(s, "O protocolo IPC",
              "Um formato de texto simples é o que torna as camadas intercambiáveis")

    mono(s, MARGEM, Inches(2.0), Inches(5.8), Inches(1.35), [
        ("casa:estado,casa:estado", TITULO),
        ("", TEXTO),
        ("e2:0,e4:1      peça saiu de e2, chegou em e4", VERDE),
    ], tam=15)

    mono(s, MARGEM, Inches(3.55), Inches(5.8), Inches(1.15), [
        ("@sync|<64 caracteres '0'/'1'>       Python → C", ACENTO),
        ("", TEXTO),
        ("Espelho de ocupação, na ordem a1..h8", FRACO),
    ], tam=15)

    cartao_texto(s, MARGEM, Inches(4.9), Inches(5.8), Inches(1.9),
                 "Eventos de exemplo", [
        "Peão e2→e4  →  e2:0,e4:1",
        "Captura Bxf7  →  c4:0,f7:1",
        "Roque curto  →  e1:0,g1:1  e depois  h1:0,f1:1",
    ], cor=VERDE)

    cartao_texto(s, MARGEM + Inches(6.2), Inches(2.0), Inches(5.6), Inches(1.35),
                 "Reed switches  (--input reed)",
                 ["O jogador move a peça no tabuleiro instrumentado."],
                 cor=VERMELHO)
    cartao_texto(s, MARGEM + Inches(6.2), Inches(3.55), Inches(5.6), Inches(3.25),
                 "Teclado matricial  (--input keypad)", [
        "O jogador digita o lance: AA 2 AA 4 #  →  e2:0,e4:1",
        "Tecla repetida vale pela letra do bloco seguinte (A→a, AA→e).",
        "Linhas @entry avisam a GUI do que está sendo digitado.",
        "",
        "Canal de volta: a aplicação manda @sync com o espelho dela, e a "
        "camada C adota — peças capturadas pelo oponente saem sozinhas.",
    ], cor=ACENTO)

    rodape(s, "A camada Python não sabe qual das duas está do outro lado — "
              "trocar é trocar uma opção de linha de comando.")


def s08_fluxo(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Fluxo de uma jogada", "Do sensor ao bestmove da engine")
    figura(s, "seq_stockfish.png", Inches(1.72), altura=Inches(4.95))
    rodape(s, "Modo Lichess: o mesmo caminho, trocando a engine local pelo "
              "POST na Board API e pelo stream de eventos da partida.")


def s09_plano_b(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Plano B: a camada de teclado",
              "Decisão de risco tomada na semana 3")

    cartao_texto(s, MARGEM, Inches(1.95), Inches(5.85), Inches(2.2),
                 "O risco identificado", [
        "A chapa de acrílico comprada ficou espessa demais para o alcance dos "
        "ímãs, e a remontagem do tabuleiro consome tempo.",
        "Avaliação do grupo: as 64 casas não fecham dentro do cronograma.",
    ], cor=VERMELHO)

    cartao_texto(s, MARGEM + Inches(6.25), Inches(1.95), Inches(5.85), Inches(2.2),
                 "A contingência", [
        "Um teclado matricial 4×4 de baixo custo em que o lance é digitado.",
        "Emite exatamente os mesmos eventos IPC — Stockfish, Lichess, GUI, "
        "roque e capturas continuam funcionando sem alteração.",
    ], cor=VERDE)

    tf = caixa(s, MARGEM, Inches(4.45), L - 2 * MARGEM, Inches(2.2))
    par(tf, "Impacto declarado sobre os requisitos", tam=19, cor=TITULO,
        negrito=True, espaco=10, primeiro=True)
    par(tf, "RF1 deixa de ser satisfeito pelo teclado — o lance é informado, "
            "não detectado — e RNF3 perde o objeto, já que não há sensor a "
            "perturbar.", tam=16, cor=TEXTO, espaco=6, marcador="•")
    par(tf, "RF2, RF3 e RNF2 continuam valendo integralmente: são exatamente "
            "eles que o teclado preserva.", tam=16, cor=TEXTO, espaco=6,
        marcador="•")
    par(tf, "A camada de reed switches continua implementada e é o caminho "
            "principal. O teclado é contingência, não troca de escopo.",
        tam=16, cor=TEXTO, espaco=0, marcador="•")


def s10_interface(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Interface com o usuário",
              "Resposta direta à revisão por pares da semana 2")

    cartao_texto(s, MARGEM, Inches(1.95), Inches(3.75), Inches(2.5),
                 "Menu de configuração", [
        "Modo de jogo, oponente, cor, camada de entrada, tempo da engine e "
        "controle de tempo escolhidos na própria interface.",
        "Janela pygame ou lista numerada no terminal (acesso por SSH).",
    ])
    cartao_texto(s, MARGEM + Inches(4.15), Inches(1.95), Inches(3.75), Inches(2.5),
                 "Sinalização de erros", [
        "Pré-condições verificadas antes de começar: token ausente, controle "
        "de tempo recusado pela Board API, processo em C não compilado, "
        "Stockfish fora do PATH.",
        "Cada uma vira mensagem no rodapé — não um erro no meio da partida.",
    ])
    cartao_texto(s, MARGEM + Inches(8.3), Inches(1.95), Inches(3.75), Inches(2.5),
                 "Ações durante a partida", [
        "Esc abre: reiniciar, abortar, oferecer ou aceitar empate, desistir, "
        "voltar ao menu, sair.",
        "Ação irreversível pede confirmação — num tabuleiro operado por toque, "
        "um encostão não pode custar a partida.",
    ])

    tf = caixa(s, MARGEM, Inches(4.85), L - 2 * MARGEM, Inches(1.9))
    par(tf, "Assistência no tabuleiro", tam=19, cor=TITULO, negrito=True,
        espaco=10, primeiro=True)
    par(tf, "Com a peça na mão, a GUI destaca os destinos legais (ponto para "
            "casa livre, anel vermelho para captura). A barra de status diz "
            "sempre o que fazer no tabuleiro físico: “Remova a peça de d5”, "
            "“Desfaça o movimento ilegal — mova a peça de e5 para e2”, "
            "“Roque — agora mova a torre de h1 para f1”.",
        tam=16, cor=TEXTO, espaco=0)


def s11_captura_auto(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Retirada automática de peça capturada",
              "Eliminando o passo manual que travava o teclado")

    larg = Inches(5.85)
    cartao_texto(s, MARGEM, Inches(1.95), larg, Inches(2.3),
                 "O problema", [
        "Quando o oponente captura uma peça do jogador, a peça sai do "
        "tabuleiro virtual mas fica no espelho da camada C.",
        "Até aqui, o jogador precisava digitar  0 1 <casa> #  para corrigir "
        "o espelho antes de poder jogar — inclusive antes de recapturar.",
    ], cor=VERMELHO)

    cartao_texto(s, MARGEM + larg + Inches(0.5), Inches(1.95), larg, Inches(2.3),
                 "A solução", [
        "Canal de volta Python → C por stdin: a aplicação manda "
        "@sync|<64 chars '0'/'1'> com o espelho dela.",
        "A camada C adota o espelho, e o LCD mostra qual peça retirar da "
        "mesa — sem bloquear o jogo.",
    ], cor=VERDE)

    tf = caixa(s, MARGEM, Inches(4.55), L - 2 * MARGEM, Inches(2.3))
    par(tf, "Onde a sincronização acontece", tam=19, cor=TITULO,
        negrito=True, espaco=10, primeiro=True)
    for item in (
        "Após todo lance aplicado (_push_mirror_to_hardware): mantém a "
        "invariante de que os dois espelhos são iguais.",
        "Após lance do oponente (_absorb_captured_pieces): esvazia do "
        "espelho Python as casas capturadas e reenvia ao C.",
        "Na reinicialização e troca de cor: o processo C nasceu com um "
        "espelho que pode não ser mais o certo.",
        "Só no teclado GPIO: reed switches leem o tabuleiro de verdade, "
        "e o mock tem tabuleiro próprio.",
    ):
        par(tf, item, tam=14, cor=TEXTO, espaco=5, marcador="▸")


def s12_demo(prs):
    s = slide_vazio(prs)
    faixa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), A)
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = ACENTO
    faixa.line.fill.background()
    faixa.shadow.inherit = False

    tf = caixa(s, Inches(1.1), Inches(2.3), Inches(11.0), Inches(3.0))
    par(tf, "DEMONSTRAÇÃO", tam=16, cor=ACENTO, negrito=True, espaco=16,
        primeiro=True)
    par(tf, "O funcionamento será mostrado ao vivo", tam=40, cor=TITULO,
        negrito=True, espaco=24)
    par(tf, "Partida contra o Stockfish pelo teclado matricial  ·  "
            "partida online via Lichess  ·  mock gráfico dos reed switches  ·  "
            "fileiras montadas do tabuleiro", tam=17, cor=TEXTO, espaco=0)


def s13_testes(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Testes realizados",
              "Cinco suítes automatizadas, sem depender de token, rede, "
              "Stockfish ou display")

    larg = Inches(5.85)
    cartao_texto(s, MARGEM, Inches(1.95), larg, Inches(3.15),
                 "Suítes  (make test)", [
        "test_keypad_layer.py — 30 verificações que executam o binário em C "
        "de verdade, digitam teclas em stdin e conferem os eventos",
        "test_keypad_capture.py — retirada automática da peça capturada: "
        "canal @sync, validação de payload, integração com a aplicação",
        "test_lichess.py e test_challenge.py — modo online contra um servidor "
        "falso da Board API",
        "test_stockfish_loop.py — regressão do loop principal com engine falsa",
        "Resultado: todas as 5 suítes passaram",
    ], cor=VERDE)

    cartao_texto(s, MARGEM + larg + Inches(0.5), Inches(1.95), larg, Inches(3.15),
                 "O que as evidências mostram", [
        "O roque funciona pelo teclado: cinco lances produzem os mesmos "
        "eventos que o tabuleiro físico produziria",
        "A troca de camada é transparente: AA2AA4# gera e2:0,e4:1, idêntico "
        "ao que a matriz de reed switches emitiria",
        "Peça capturada pelo oponente sai do espelho sozinha, sem intervenção "
        "manual — e o LCD mostra qual peça tirar da mesa",
        "Canal @sync não gera evento de volta — a extensão é retrocompatível",
        "Payloads malformados (curto, inválido, desconhecido) são ignorados: "
        "meio espelho aplicado nunca acontece",
    ])

    cartao_texto(s, MARGEM, Inches(5.35), L - 2 * MARGEM, Inches(1.5),
                 "Validação de hardware", [
        "Teclado matricial testado com hardware real (linhas BCM 16/20/21/26, "
        "colunas 19/13/6/5), conferido com o modo de bancada --raw, que mostra "
        "em que interseção da matriz cada tecla foi lida.",
        "Tabuleiro instrumentado: cerca de duas fileiras completas montadas com "
        "reed switches e diodos — suficiente para confirmar em escala real o "
        "mapeamento de linhas e colunas e a detecção através da chapa.",
    ], cor=ACENTO)


def s14_rastreabilidade(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Resultados por requisito", "Rastreabilidade requisito × evidência")

    linhas = [
        ("RF1", "Detecção das peças e movimentos",
         "Software validado (mock + loop principal); hardware parcial — 2 fileiras montadas",
         ACENTO),
        ("RF2", "Comunicação com a chess engine",
         "Atendido — Stockfish via UCI e Lichess Board API, cobertos por testes",
         VERDE),
        ("RF3", "Exibição gráfica das duas partes",
         "Atendido — GUI pygame com 32 peças, destaques e barra de instruções",
         VERDE),
        ("RNF1", "Detecção < 200 ms e 100 % de precisão",
         "Precisão 30/30 nos eventos verificados; medição de latência pendente",
         ACENTO),
        ("RNF2", "GUI sem latência considerável",
         "Renderização incremental implementada; instrumentação de tempos pendente",
         ACENTO),
        ("RNF3", "Robustez contra falsos positivos",
         "Dupla barreira (debouncing em C + validação em Python), ressincronização "
         "e espelho bidirecional (canal @sync); debouncing em hardware pendente", ACENTO),
    ]
    y = Inches(1.95)
    alt = Inches(0.62)
    for req, nome, estado, cor in linhas:
        f = cartao(s, MARGEM, y, L - 2 * MARGEM, alt, cor_borda=BORDA)
        tf = f.text_frame
        tf.margin_left = Inches(0.22)
        tf.margin_top = tf.margin_bottom = Inches(0.05)
        tf.vertical_anchor = MSO_ANCHOR.MIDDLE
        p = tf.paragraphs[0]
        p.space_after = Pt(0)
        for texto, tam, c, negrito in (
                (f"{req}   ", 15, cor, True),
                (f"{nome}      ", 14, TITULO, False),
                (estado, 13, TEXTO, False)):
            r = p.add_run()
            r.text = texto
            r.font.size = Pt(tam)
            r.font.bold = negrito
            r.font.color.rgb = c
        y += alt + Inches(0.10)

    rodape(s, "As pendências são de medição, não de implementação: exigem "
              "instrumentar o caminho completo no Raspberry Pi.")


def s15_dificuldades(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Dificuldades encontradas", "O que custou tempo e o que aprendemos")

    itens = [
        ("Mecânica antes da eletrônica",
         "A chapa de acrílico escolhida ficou espessa demais para o alcance dos "
         "ímãs de neodímio. Um detalhe mecânico — não de software nem de "
         "circuito — foi o que travou o cronograma do hardware."),
        ("Montagem de 64 casas é trabalho manual",
         "Cada casa exige reed switch, diodo e solda. O esforço cresce linearmente "
         "e não paraleliza bem entre os integrantes."),
        ("Tabuleiro físico e tabuleiro virtual divergem",
         "Reed switches dizem onde há ímã, não qual peça é. Capturas do oponente "
         "e lances ilegais dessincronizam os dois — resolvido com histórico de "
         "peças deslocadas, instruções na barra de status, ressincronização e "
         "canal @sync bidirecional para a camada de teclado."),
        ("Restrições da Board API",
         "O Lichess recusa controles de tempo abaixo de 480 s estimados. "
         "Descoberto empiricamente e transformado em verificação prévia, com "
         "explicação, em vez de um HTTP 400 sem contexto."),
        ("Testar hardware que ainda não existe",
         "Resolvido com mock gráfico dos sensores, stubs de GPIO que permitem "
         "compilar o binário fora do Raspberry Pi e um servidor falso da Board API."),
    ]
    y = Inches(1.9)
    for titulo, corpo in itens:
        m = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, MARGEM, y + Inches(0.06),
                               Pt(3.5), Inches(0.72))
        m.fill.solid()
        m.fill.fore_color.rgb = ACENTO
        m.line.fill.background()
        m.shadow.inherit = False
        tf = caixa(s, MARGEM + Inches(0.28), y, L - 2 * MARGEM - Inches(0.28),
                   Inches(0.9))
        par(tf, titulo, tam=17, cor=TITULO, negrito=True, espaco=2, primeiro=True)
        par(tf, corpo, tam=13, cor=TEXTO, espaco=0)
        y += Inches(1.05)


def s16_futuro(prs):
    s = slide_vazio(prs)
    cabecalho(s, "Trabalhos futuros", "O que fica encaminhado")

    larg = Inches(3.75)
    cartao_texto(s, MARGEM, Inches(1.95), larg, Inches(3.1),
                 "Hardware", [
        "Concluir as 64 casas do tabuleiro com chapa de espessura adequada",
        "Medir a taxa de acerto da detecção e exercitar o debouncing com "
        "esbarrões reais (RNF3)",
        "Integrar o display 16×2 por I²C ao caminho principal",
    ], cor=VERMELHO)
    cartao_texto(s, MARGEM + larg + Inches(0.4), Inches(1.95), larg, Inches(3.1),
                 "Medição", [
        "Instrumentar o caminho completo — sensor ou tecla até o quadro "
        "renderizado — no próprio Raspberry Pi",
        "Fechar as evidências de RNF1 (< 200 ms) e RNF2 (< 100 ms)",
        "Teste de aceitação: partidas completas até o xeque-mate",
    ], cor=ACENTO)
    cartao_texto(s, MARGEM + 2 * (larg + Inches(0.4)), Inches(1.95), larg,
                 Inches(3.1), "Software", [
        "Envio de lances ao Lichess fora do loop principal (hoje é POST síncrono)",
        "Reconexão automática do stream e reconstrução do estado após takeback",
        "Diálogo de promoção (hoje promove sempre para dama)",
        "Relógios da partida e suporte a variantes",
    ], cor=VERDE)

    tf = caixa(s, MARGEM, Inches(5.4), L - 2 * MARGEM, Inches(1.4))
    par(tf, "Considerações finais", tam=19, cor=TITULO, negrito=True, espaco=8,
        primeiro=True)
    par(tf, "O sistema está jogável de ponta a ponta — contra o Stockfish e "
            "contra oponentes online — com toda a lógica de xadrez, a interface "
            "e a integração validadas por testes automatizados. A retirada "
            "automática de peças capturadas fechou o último atrito de usabilidade "
            "do teclado. O que separa o projeto do escopo original é a montagem "
            "física do tabuleiro, e a decisão de manter as camadas de entrada "
            "intercambiáveis foi o que permitiu que esse atraso não bloqueasse "
            "o resto.",
        tam=15, cor=TEXTO, espaco=0)


def s17_fim(prs):
    s = slide_vazio(prs)
    faixa = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, Inches(0.22), A)
    faixa.fill.solid()
    faixa.fill.fore_color.rgb = ACENTO
    faixa.line.fill.background()
    faixa.shadow.inherit = False

    tf = caixa(s, Inches(1.1), Inches(2.6), Inches(11.0), Inches(2.4))
    par(tf, "Obrigado", tam=48, cor=TITULO, negrito=True, espaco=18,
        primeiro=True)
    par(tf, "Perguntas?", tam=24, cor=ACENTO, espaco=24)
    par(tf, "Grupo W  ·  PCS3732 — Laboratório de Processadores  ·  "
            "Escola Politécnica da USP", tam=15, cor=FRACO, espaco=0)


def main():
    prs = Presentation()
    prs.slide_width, prs.slide_height = L, A
    for construir in (s01_capa, s02_motivacao, s03_objetivos, s04_solucao,
                      s05_arquitetura, s06_fisica, s07_ipc, s08_fluxo,
                      s09_plano_b, s10_interface, s11_captura_auto,
                      s12_demo, s13_testes, s14_rastreabilidade,
                      s15_dificuldades, s16_futuro, s17_fim):
        construir(prs)
    prs.save(SAIDA)
    print(f"{SAIDA.relative_to(RAIZ)} — {len(prs.slides.__iter__.__self__._sldIdLst)} slides")


if __name__ == "__main__":
    main()
